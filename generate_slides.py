"""
Generate slide deck from Webinar_Template.pptx
- Slide 1 (title) and Slide 2 (agenda): kept as-is
- Slide 3 layout (CUSTOM_5_1_1): section divider with large centered text
- Slides 4/5 layout (CUSTOM_2_2_2 / CUSTOM_2_2_2_3): content slide
    Left: main text area (ph_idx=0)
    Right: 4 label/description pairs (ph_idx 2-9)
- Original slides 3-6 are removed; new slides appended after slide 2
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn


def delete_slide(prs, slide):
    """Delete a slide from the presentation."""
    rId = None
    for rel in prs.part.rels.values():
        if rel.target_part == slide.part:
            rId = rel.rId
            break
    if rId:
        sldIdLst = prs.part._element.find(qn('p:sldIdLst'))
        for sldId in sldIdLst.findall(qn('p:sldId')):
            if sldId.get(qn('r:id')) == rId:
                sldIdLst.remove(sldId)
                break


def set_placeholder_text(slide, ph_idx, text):
    """Fill a placeholder with text."""
    try:
        ph = slide.placeholders[ph_idx]
    except KeyError:
        return
    ph.text = text


def add_section_slide(prs, layout, title_text):
    """Add a section divider slide (slide 3 style)."""
    slide = prs.slides.add_slide(layout)
    set_placeholder_text(slide, 0, title_text)
    return slide


def add_content_slide(prs, layout, left_text, pairs, notes=""):
    """
    Add a content slide (slide 4/5 style).
    left_text: main text for left area (ph_idx=0)
    pairs: list of (label, description) tuples, up to 4
    notes: speaker notes text
    """
    slide = prs.slides.add_slide(layout)
    set_placeholder_text(slide, 0, left_text)

    ph_map = [(2, 3), (4, 5), (6, 7), (8, 9)]
    for i, (label, desc) in enumerate(pairs[:4]):
        if i < len(ph_map):
            label_idx, desc_idx = ph_map[i]
            set_placeholder_text(slide, label_idx, label)
            set_placeholder_text(slide, desc_idx, desc)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


def main():
    prs = Presentation("Webinar_Template.pptx")

    # Capture layout references from existing slides before deleting
    section_layout = prs.slides[2].slide_layout   # Slide 3: CUSTOM_5_1_1
    content_layout_a = prs.slides[3].slide_layout  # Slide 4: CUSTOM_2_2_2
    content_layout_b = prs.slides[4].slide_layout  # Slide 5: CUSTOM_2_2_2_3

    # Delete original slides 3-6 (reverse order to avoid index shifts)
    for i in [5, 4, 3, 2]:
        delete_slide(prs, prs.slides[i])

    # ── SECTION: What Is an AI Knowledge Graph? ───────────────

    add_section_slide(prs, section_layout,
                      "What Is an AI\nKnowledge Graph?")

    add_content_slide(prs, content_layout_a,
        left_text=(
            "A knowledge graph organizes information "
            "as entities (nodes) and relationships (edges), "
            "stored as JSON documents in MongoDB."
        ),
        pairs=[
            ("Entities (Nodes)",
             "Documents, people, concepts, companies — any object with properties"),
            ("Relationships (Edges)",
             "Cross-references, resource links, dependencies between entities"),
            ("Properties & Metadata",
             "Attributes on both nodes and edges (type, date, service line, etc.)"),
            ("JSON Representation",
             '{"entity": "MongoDB", "relations": [{"type": "PRODUCES", "target": "Atlas"}]}'),
        ],
        notes=(
            "A knowledge graph organizes information as entities and relationships. "
            "In MongoDB these are simply JSON documents. Each entity can embed arrays "
            "of relationships pointing to other entities. Example: Elon Musk with "
            "relations FOUNDED → SpaceX and CEO_OF → Tesla."
        ))

    add_content_slide(prs, content_layout_b,
        left_text=(
            "AI adds semantic understanding.\n"
            "Graphs add structure.\n"
            "Together they find what you\n"
            "didn't know to search for."
        ),
        pairs=[
            ("Traditional Search",
             "Keyword matching only, requires exact terms, no context awareness"),
            ("AI Knowledge Graph",
             "Semantic understanding, multi-hop reasoning, relationship discovery"),
            ("Key Advantage",
             "Finds conceptually related content even with different terminology"),
            ("Result",
             "Explainable, traceable results with full provenance across hops"),
        ],
        notes=(
            "Traditional keyword search only finds exact matches. An AI knowledge graph "
            "combines semantic understanding from vector embeddings with structural "
            "relationships from the graph. You discover documents you didn't know existed."
        ))

    # ── SECTION: MongoDB as a Unified Platform ────────────────

    add_section_slide(prs, section_layout,
                      "MongoDB as a\nUnified Platform")

    add_content_slide(prs, content_layout_a,
        left_text=(
            "MongoDB Atlas unifies operational "
            "data, vector embeddings, and graph "
            "relationships in a single platform.\n\n"
            "No separate vector DB.\n"
            "No separate graph DB."
        ),
        pairs=[
            ("Operational Data",
             "JSON documents with full CRUD and aggregation framework"),
            ("Vector Embeddings",
             "Atlas Vector Search with cosine similarity for semantic search"),
            ("Graph Relationships",
             "$graphLookup for recursive multi-hop traversal up to N depth"),
            ("Unified Platform",
             "One database for all three — zero data duplication across systems"),
        ],
        notes=(
            "MongoDB Atlas serves as a unified platform. Operational data is standard "
            "JSON with full CRUD. Vector embeddings power semantic search. Graph "
            "relationships are traversed using $graphLookup. No extra infrastructure needed."
        ))

    add_content_slide(prs, content_layout_b,
        left_text=(
            "Two data modeling approaches\n"
            "for knowledge graphs in MongoDB.\n\n"
            "The demo app supports both —\n"
            "toggle between them at search time."
        ),
        pairs=[
            ("Approach A: Document-Centric",
             "Embedded connections array within each document. Best for sparse graphs."),
            ("Approach B: Graph Model",
             "Separate graph_nodes + graph_edges collections. Best for dense graphs."),
            ("Trade-offs",
             "A is simpler/faster for reads. B scales better for complex traversals."),
            ("Flexibility",
             "Same search and agent features work across both models seamlessly."),
        ],
        notes=(
            "Two data modeling approaches. Approach A embeds connections directly "
            "in each document. Approach B normalizes into graph_nodes and graph_edges. "
            "The demo app lets you switch between them."
        ))

    add_content_slide(prs, content_layout_a,
        left_text=(
            "GraphRAG combines vector search\n"
            "with graph traversal to discover\n"
            "documents the user didn't know\n"
            "to search for.\n\n"
            "Improves accuracy, explainability,\n"
            "and multi-hop reasoning."
        ),
        pairs=[
            ("1. Embed the Query",
             "Convert user question into a vector using VoyageAI or OpenAI"),
            ("2. Vector Search",
             "Find semantically similar seed documents via Atlas Vector Search"),
            ("3. Graph Expansion",
             "Traverse relationships with $graphLookup to discover connected docs"),
            ("4. Combined Context",
             "Merge vector + graph results for rich, multi-hop context"),
        ],
        notes=(
            "GraphRAG: embed the query, vector search for seeds, $graphLookup to expand, "
            "combined context. This finds documents the user didn't know to search for."
        ))

    add_content_slide(prs, content_layout_b,
        left_text=(
            "Full-stack application built for\n"
            "enterprise knowledge graph\n"
            "management and AI-powered\n"
            "document analysis."
        ),
        pairs=[
            ("Frontend",
             "React 19 + D3.js (graph viz) + Tailwind CSS"),
            ("Backend",
             "Express 5 / Node.js with RESTful API"),
            ("Database & Search",
             "MongoDB Atlas · Atlas Search · Vector Search · $rankFusion"),
            ("AI / Embeddings",
             "VoyageAI (voyage-4-lite, 1024d) or OpenAI (ada-002, 1536d)"),
        ],
        notes=(
            "Tech stack: React 19 with D3.js for visualization, Express 5 backend, "
            "MongoDB Atlas with Atlas Search, Vector Search, and $graphLookup. "
            "Supports VoyageAI and OpenAI for embeddings."
        ))

    # ── SECTION: Demo — Connections Graph ─────────────────────

    add_section_slide(prs, section_layout,
                      "Demo: Connections Graph")

    # -- Setup & Indexes --

    add_content_slide(prs, content_layout_a,
        left_text=(
            "Setup: Building the\n"
            "Knowledge Graph\n\n"
            "Sequential pipeline from\n"
            "connection to search-ready\n"
            "knowledge graph."
        ),
        pairs=[
            ("1. Connect",
             "Provide MongoDB Atlas URI and select database name"),
            ("2. Seed Data",
             "Generate 120 synthetic documents with rich metadata"),
            ("3. Create Indexes",
             "Database indexes + Atlas Search + Vector Search indexes"),
            ("4. Generate Embeddings",
             "Batch embed all docs (32/batch) via VoyageAI or OpenAI"),
        ],
        notes=(
            "LIVE DEMO: Walk through Setup Page. Connect → Seed 120 docs → "
            "Create 3 types of indexes → Generate embeddings in batches of 32. "
            "Status dashboard shows real-time progress."
        ))

    add_content_slide(prs, content_layout_b,
        left_text=(
            "Under the Hood: Indexing\n\n"
            "Both index types live on the\n"
            "same collection — zero data\n"
            "duplication required."
        ),
        pairs=[
            ("Atlas Search Index",
             "Full-text on title, topic, category, serviceline, focus. Fuzzy matching."),
            ("Vector Search Index",
             "Cosine similarity on embeddings. 1024d (VoyageAI) or 1536d (OpenAI)."),
            ("Compound Queries",
             "Combine text filters with vector scoring in a single pipeline."),
            ("Pipeline Transparency",
             "The app shows the exact aggregation pipeline for every search."),
        ],
        notes=(
            "Atlas Search indexes enable full-text with fuzzy matching. Vector Search "
            "indexes enable semantic similarity using cosine distance. Both are on the "
            "same collection. The app shows the pipeline for transparency."
        ))

    # -- Four Search Strategies --

    add_content_slide(prs, content_layout_a,
        left_text=(
            "Four Ways to Search\n\n"
            "Progressive complexity from\n"
            "simple keyword matching to\n"
            "full GraphRAG discovery.\n\n"
            "Run the same query through\n"
            "all four — compare results."
        ),
        pairs=[
            ("1. Lexical Search",
             "Atlas Search with fuzzy keyword matching for exact term lookups"),
            ("2. Vector Search",
             "Semantic similarity via embeddings — finds related concepts"),
            ("3. Hybrid ($rankFusion)",
             "Combines text + vector results 50/50 server-side — native operator"),
            ("4. Hybrid Graph",
             "Vector seeds → $graphLookup expansion for multi-hop discovery"),
        ],
        notes=(
            "LIVE DEMO: Run the same query through all four strategies. "
            "Lexical finds keywords. Vector finds similar concepts. "
            "Hybrid merges both with $rankFusion. Hybrid Graph adds graph expansion."
        ))

    add_content_slide(prs, content_layout_b,
        left_text=(
            "Hybrid Search with $rankFusion\n\n"
            "Native MongoDB operator merges\n"
            "text and vector results server-side.\n"
            "No client-side post-processing."
        ),
        pairs=[
            ("Text Results",
             "Ranked by Atlas Search relevance score (keyword matching)"),
            ("Vector Results",
             "Ranked by cosine similarity score (semantic meaning)"),
            ("$rankFusion (50/50)",
             "Combines both ranked lists into unified results — configurable weights"),
            ("Pipeline Viewer",
             "See the exact MongoDB aggregation pipeline executed for every query"),
        ],
        notes=(
            "$rankFusion is a native MongoDB operator — no client-side merging. "
            "Show the pipeline viewer to demonstrate the actual aggregation pipeline."
        ))

    # -- Graph Visualization --

    add_content_slide(prs, content_layout_a,
        left_text=(
            "Interactive Knowledge Graph\n\n"
            "D3.js powered visualization\n"
            "with multiple layouts and\n"
            "click-through exploration.\n\n"
            "Nodes sized by degree.\n"
            "Edges styled by link type."
        ),
        pairs=[
            ("Force-Directed Layout",
             "Physics-based simulation — nodes repel, edges attract"),
            ("Hierarchical Tree",
             "Parent-child structure showing document dependencies"),
            ("Radial Circle",
             "Circular arrangement for pattern recognition at a glance"),
            ("Interactions",
             "Click nodes to load neighbors, zoom/pan, inspect full metadata"),
        ],
        notes=(
            "LIVE DEMO: Show D3.js graph visualization. Nodes colored by collection, "
            "sized by degree. Edge styles: solid = Cross Reference, dashed = Image Link, "
            "dotted = Resource Link. Toggle layouts. Click-through exploration."
        ))

    # -- Agents --

    add_content_slide(prs, content_layout_b,
        left_text=(
            "Seven Intelligent Agents\n\n"
            "Specialized analysis tools\n"
            "that go beyond search to surface\n"
            "insights, risks, and structure\n"
            "from the knowledge graph."
        ),
        pairs=[
            ("Natural Language Query",
             "Ask questions in plain English — the agent interprets and searches"),
            ("Document Impact & Cross-Ref",
             "Blast radius analysis, hub detection, orphan nodes, bidirectional links"),
            ("Similarity & Collection",
             "Find semantically similar docs. Browse collection relationships."),
            ("Critical Nodes & Stale Docs ★",
             "Identify structural backbones and compliance risks (demo next)"),
        ],
        notes=(
            "Seven agents for specialized analysis. We'll demo three: "
            "Document Impact, Critical Node Detection, and Stale Document Detector."
        ))

    add_content_slide(prs, content_layout_a,
        left_text=(
            "Document Impact Analysis\n\n"
            "\"If this policy changes, which\n"
            "downstream documents are\n"
            "affected?\"\n\n"
            "Maps the full blast radius of\n"
            "any document change."
        ),
        pairs=[
            ("Outgoing Connections",
             "Documents that THIS document references directly"),
            ("Incoming Connections",
             "Documents that REFERENCE this document — downstream impact"),
            ("Grouped by Link Type",
             "Cross-references, resource links, image links broken out separately"),
            ("Change Management",
             "Know the full impact before updating any policy or standard"),
        ],
        notes=(
            "LIVE DEMO: Pick a document. The Impact agent shows outgoing + incoming "
            "edges grouped by link type. Critical for change management."
        ))

    add_content_slide(prs, content_layout_b,
        left_text=(
            "Critical Node Detection\n\n"
            "Identifies structurally important\n"
            "documents using degree centrality\n"
            "and bridge scoring.\n\n"
            "Score = totalDegree × 1\n"
            "          + bridgeCount × 2\n"
            "          + linkTypeCount × 0.5"
        ),
        pairs=[
            ("Total Degree (× 1)",
             "Number of inbound + outbound connections — raw connectivity"),
            ("Bridge Count (× 2)",
             "Connects otherwise isolated clusters — highest weight factor"),
            ("Link Type Count (× 0.5)",
             "Diversity of relationship types indicates broad importance"),
            ("Result",
             "Ranked list of backbone documents — protect and review these first"),
        ],
        notes=(
            "LIVE DEMO: Critical Node Detection ranks documents by structural importance. "
            "Bridge nodes weighted 2x because losing them fragments the network."
        ))

    add_content_slide(prs, content_layout_a,
        left_text=(
            "Stale Document Detector\n\n"
            "Compliance risk: finds old\n"
            "documents still heavily\n"
            "referenced by active content.\n\n"
            "Risk Score = inboundRefs\n"
            "  × (1 + daysPastCutoff / 365)"
        ),
        pairs=[
            ("The Risk",
             "A 3-year-old policy is still referenced by 47 active documents"),
            ("Risk Scoring",
             "inboundRefs × (1 + daysPastCutoff/365) — higher = more urgent"),
            ("Example",
             "47 refs × (1 + 730/365) = 141 — flagged as HIGH RISK"),
            ("Compliance Use Case",
             "Regulatory compliance, audit trail, proactive risk management"),
        ],
        notes=(
            "LIVE DEMO: Stale Document Detector is the compliance hook. "
            "Finds docs past review date still referenced heavily. "
            "Directly relevant to audit and regulatory compliance."
        ))

    # ── SECTION: Closing ──────────────────────────────────────

    add_section_slide(prs, section_layout,
                      "Key Takeaways & Q&A")

    add_content_slide(prs, content_layout_b,
        left_text=(
            "One platform.\n"
            "Three capabilities.\n"
            "Zero silos.\n\n"
            "Thank you!\n"
            "Q&A"
        ),
        pairs=[
            ("Unified Platform",
             "MongoDB Atlas combines docs, vectors, and graph — no extra infrastructure"),
            ("Native Operators",
             "$rankFusion + $graphLookup run inside the DB, not in application code"),
            ("AI-Powered Agents",
             "Surface compliance risks, critical nodes, and impact automatically"),
            ("Next Steps",
             "Contact info · Resources · Demo repository access"),
        ],
        notes=(
            "Recap: MongoDB Atlas eliminates separate vector and graph databases. "
            "Native operators keep computation in the DB. AI agents automate analysis. "
            "One platform, three capabilities, zero silos. Open for questions."
        ))

    # ── Save ──────────────────────────────────────────────────
    output = "AI_Knowledge_Graphs_MongoDB.pptx"
    prs.save(output)
    print(f"✅ Saved {output} — {len(prs.slides)} slides")
    for i, slide in enumerate(prs.slides):
        filled = sum(1 for ph in slide.placeholders if ph.text.strip())
        ln = slide.slide_layout.name
        print(f"  Slide {i+1}: layout={ln[:35]:<35} filled_placeholders={filled}")


if __name__ == "__main__":
    main()
